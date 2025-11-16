import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import ollama
import json
import re
from pathlib import Path
from config import OLLAMA_API_KEY, OLLAMA_MODEL

class PresentationGenerator:
    def __init__(self, templates_dir="./presentations"):
        self.client = ollama.Client(
            host="https://ollama.com",
            headers={'Authorization': f'Bearer {OLLAMA_API_KEY}'}
        )
        self.templates_dir = Path(templates_dir)
        if not self.templates_dir.exists():
            self.templates_dir.mkdir()
    
    def get_available_templates(self):
        """Возвращает список доступных шаблонов"""
        templates = []
        for file in self.templates_dir.glob("*.pptx"):
            templates.append({
                "name": file.stem,
                "path": file
            })
        return templates
    
    def extract_template_info(self, template_path):
        """Извлекает структуру и стиль из шаблона"""
        try:
            prs = Presentation(template_path)
            
            # Извлекаем структуру
            structure = {"slides": []}
            for slide in prs.slides:
                slide_info = {
                    "title": slide.shapes.title.text if slide.shapes.title else "",
                    "points": []
                }
                
                # Ищем контент
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text and len(text) > 10:
                                slide_info["points"].append(text)
                
                structure["slides"].append(slide_info)
            
            # Определяем стиль
            style = "business"
            if "creative" in template_path.stem.lower():
                style = "creative"
            elif "minimal" in template_path.stem.lower():
                style = "minimal"
            elif "dark" in template_path.stem.lower():
                style = "dark"
            
            return structure, style
            
        except Exception as e:
            print(f"Ошибка чтения шаблона {template_path}: {e}")
            return None, "business"
    
    def generate_structure(self, topic, language="ru", slides_count=5, template_structure=None):
        """Генерирует структуру презентации через Ollama"""
        lang_prompt = "на русском языке" if language == "ru" else "in English"
        
        template_info = ""
        if template_structure:
            template_info = f"""
ВАЖНО: Используй структуру этого шаблона как основу:
{json.dumps(template_structure, ensure_ascii=False, indent=2)}"""
        
        prompt = f"""Создай подробный план презентации {lang_prompt} на тему "{topic}".
Количество слайдов: {slides_count}.{template_info}

Для каждого слайда укажи: заголовок и 3-5 ключевых пунктов.
Отвечай СТРОГО в формате JSON:
{{"slides": [{{"title": "...", "points": ["...", "..."]}}]}}"""
        
        try:
            response = self.client.chat(
                model=OLLAMA_MODEL,
                messages=[{'role': 'user', 'content': prompt}],
                stream=False
            )
            
            content = response['message']['content']
            print(f"Ответ модели: {content[:200]}...")
            
            json_match = re.search(r'\{.*\}', content, re.DOTALL)
            if json_match:
                data = json.loads(json_match.group())
                if self._validate_structure(data, slides_count):
                    return data
                    
        except Exception as e:
            print(f"Ошибка генерации: {e}")
        
        return self._fallback_structure(topic, slides_count)
    
    def _validate_structure(self, data, expected_slides):
        """Проверяет корректность структуры"""
        if not isinstance(data, dict) or "slides" not in data:
            return False
        if len(data["slides"]) != expected_slides:
            return False
        for slide in data["slides"]:
            if "title" not in slide:
                return False
        return True
    
    def _fallback_structure(self, topic, slides_count):
        """Создаёт аварийную структуру"""
        return {
            "slides": [
                {"title": topic, "subtitle": "AI Generated"},
                *[{
                    "title": f"Аспект {i}",
                    "points": ["Описание проблемы", "Решение", "Результат"]
                } for i in range(1, slides_count)]
            ]
        }
    
    def create_presentation(self, structure_data, style="business", template_path=None):
        """
        Создаёт презентацию в памяти, используя шаблон (опционально)
        Возвращает BytesIO buffer (не сохраняет на диск)
        """
        if template_path and Path(template_path).exists():
            # Открываем шаблон
            prs = Presentation(template_path)
            
            # Очищаем контент (оставляем только 1 слайд)
            while len(prs.slides) > 1:
                rId = prs.slides._sldIdLst[-1].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[-1]
            
            # Обновляем титульный слайд
            if len(prs.slides) > 0:
                title_slide = prs.slides[0]
                if title_slide.shapes.title:
                    title_slide.shapes.title.text = structure_data["slides"][0]["title"]
        else:
            # Создаём новую презентацию
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
            # Титульный слайд
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = structure_data["slides"][0]["title"]
            subtitle.text = "AI-сгенерированная презентация"
        
        # Добавляем остальные слайды
        for slide_info in structure_data["slides"][1:]:
            slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            # Заголовок
            if slide.shapes.title:
                slide.shapes.title.text = slide_info["title"]
            
            # Контент
            if len(slide.placeholders) > 1:
                content = slide.placeholders[1].text_frame
                content.clear()
                
                for point in slide_info.get("points", []):
                    p = content.add_paragraph()
                    p.text = f"• {point}"
                    p.level = 0
        
        # Сохраняем в буфер (не на диск!)
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        return pptx_buffer